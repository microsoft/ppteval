"""
Code-writing agent for Python-pptx environment.

Agent that generates Python code snippets to manipulate PowerPoint files.
Works with PythonPptxEnvironment (CodeState).
"""

import json
import logging
from typing import Any

from ppteval.core.base import Action, Agent, State
from ppteval.environments.pythonpptx_environment import CodeState


class CodeAgent(Agent):
    """
    Base agent for code-based PowerPoint manipulation.

    Subclasses should implement the code generation logic.
    This agent receives CodeState (execution results + PPT info)
    and returns Actions with Python code to execute.
    """

    def __init__(self, config: dict[str, Any] | None = None):
        """
        Initialize code agent.

        Args:
            config: Agent configuration
        """
        self.config = config or {}
        self.logger = logging.getLogger(self.__class__.__name__)
        self.instruction: str | None = None

    def step(self, state: State) -> Action:
        """
        Take a step given current state.

        Args:
            state: Current CodeState with execution results

        Returns:
            Action with Python code to execute

        Raises:
            ValueError: If state is not CodeState or instruction not set
        """
        if not isinstance(state, CodeState):
            raise ValueError(f"CodeAgent requires CodeState, got {type(state)}")

        if self.instruction is None:
            raise ValueError("Instruction not set. Call set_instruction() first.")

        # Format state information for the agent
        state_info = self._format_state(state)

        # Generate code (subclasses implement this)
        code = self._generate_code(state_info)

        # Check if agent wants to finish
        if self._should_finish(state):
            return Action(
                action_type="finish",
                params={"message": "Task completed"},
                reasoning="Generated code completed the task",
            )

        # Return code execution action
        return Action(
            action_type="execute_code",
            params={"code": code},
            reasoning=f"Executing python-pptx code to: {self.instruction[:100]}",
        )

    def _format_state(self, state: CodeState) -> str:
        """
        Format CodeState into text description for agent.

        Args:
            state: Current CodeState

        Returns:
            Formatted string with state information
        """
        parts = []

        # Execution result from last step
        if state.output:
            parts.append(f"Last execution output:\n{state.output}\n")

        if state.error:
            parts.append(f"Last execution error:\n{state.error}\n")

        # PPT structure
        parts.append("Current PowerPoint structure:")
        parts.append(f"- Slide count: {state.ppt_info.get('slide_count', 0)}")

        for slide_info in state.ppt_info.get('slides', []):
            slide_num = slide_info.get('slide_number', '?')
            shape_count = slide_info.get('shape_count', 0)
            parts.append(f"- Slide {slide_num}: {shape_count} shapes")

            for shape_info in slide_info.get('shapes', [])[:5]:  # First 5 shapes
                shape_id = shape_info.get('shape_id', '?')
                shape_type = shape_info.get('shape_type', 'unknown')
                text = shape_info.get('text', '')
                if text:
                    parts.append(f"  - Shape {shape_id} ({shape_type}): {text[:50]}...")

        return '\n'.join(parts)

    def _generate_code(self, state_info: str) -> str:
        """
        Generate Python code for next step.

        Subclasses must implement this.

        Args:
            state_info: Formatted state information

        Returns:
            Python code to execute
        """
        raise NotImplementedError("Subclasses must implement _generate_code()")

    def _should_finish(self, state: CodeState) -> bool:
        """
        Check if agent should finish.

        Subclasses can override this.

        Args:
            state: Current CodeState

        Returns:
            True if task should be considered complete
        """
        return False

    def reset(self) -> None:
        """Reset agent state for new task."""
        self.instruction = None

    def set_instruction(self, instruction: str) -> None:
        """
        Set task instruction.

        Args:
            instruction: Task description
        """
        self.instruction = instruction
        self.logger.debug(f"Set instruction: {instruction[:100]}...")

    def close(self) -> None:
        """Clean up agent resources."""
        pass


class LLMCodeAgent(CodeAgent):
    """
    LLM-based code generation agent using LiteLLM (supports any model).

    Generates python-pptx code to manipulate PowerPoint files.
    Uses LiteLLM to support Claude, GPT-4o, GPT-4.5, etc.
    """

    def __init__(self, config: dict[str, Any] | None = None):
        """
        Initialize LLM code agent with LiteLLM.

        Args:
            config: Configuration including:
                - model: Model name (e.g., "claude-sonnet-4", "gpt-4o", "gpt-4.5")
                        Supports LiteLLM format: azure/deployment, openai/model, etc.
                - temperature: Sampling temperature (default: 0.7)
                - max_tokens: Max response tokens (default: 2000)
                - api_key: API key for the provider (or use environment variables)
                - api_base: Base URL for API (optional)

        Environment variables used by LiteLLM:
            - ANTHROPIC_API_KEY: For Claude models
            - OPENAI_API_KEY: For OpenAI models
            - AZURE_API_KEY: For Azure OpenAI
            - AZURE_API_BASE: For Azure OpenAI endpoint
        """
        super().__init__(config)

        import os

        # Model configuration
        self.model = self.config.get("model", "claude-sonnet-4-20250514")
        self.temperature = self.config.get("temperature", 0.7)
        self.max_tokens = self.config.get("max_tokens", 2000)

        # Optional API configuration (LiteLLM uses env vars by default)
        self.api_key = self.config.get("api_key")
        self.api_base = self.config.get("api_base")

        self.logger.info(f"Initialized LLMCodeAgent with model: {self.model}")

        self.conversation_history = []
        self.code_attempts = 0

    def _generate_code(self, state_info: str) -> str:
        """
        Use Azure OpenAI GPT-4o to generate Python code.

        Args:
            state_info: Current state information

        Returns:
            Generated Python code
        """
        self.code_attempts += 1

        # Build system prompt with python-pptx expertise
        system_prompt = """You are an expert Python programmer specializing in the python-pptx library for PowerPoint manipulation.

Your task is to generate clean, working Python code that uses python-pptx to accomplish specific goals.

Key facts about the environment:
- The variable 'prs' is already loaded as a Presentation object
- You can access slides via prs.slides[index] or iterate with for slide in prs.slides
- Common operations:
  - Add slide: prs.slides.add_slide(layout)
  - Access shapes: slide.shapes[index] or iterate
  - Modify text: shape.text = "new text" or shape.text_frame.text = "new text"
  - Add shapes: slide.shapes.add_shape(shape_type, left, top, width, height)
  - Access placeholders: slide.placeholders[index]
- Use pptx.util.Inches() and pptx.util.Pt() for measurements
- Import what you need from pptx (e.g., from pptx.util import Inches, Pt)
- The presentation is automatically saved after your code runs
- Use print() to output status messages

Rules:
1. Write ONLY executable Python code - no markdown, no explanations, no code blocks
2. Make ONE logical change at a time to avoid errors
3. Use defensive programming (check if slides/shapes exist before accessing)
4. Print what you're doing for debugging
5. Handle potential errors gracefully
6. If previous code failed, try a different approach

Output the raw Python code that can be executed directly."""

        # Build user prompt with task and current state
        user_prompt = f"""Task to accomplish:
{self.instruction}

Current PowerPoint state:
{state_info}

Generate Python code (raw, no markdown) to make progress on this task. The 'prs' variable is already loaded."""

        # Add to conversation history
        if not self.conversation_history:
            # First message
            self.conversation_history = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ]
        else:
            # Subsequent messages - add the new state info
            self.conversation_history.append({
                "role": "user",
                "content": f"Previous code execution result:\n{state_info}\n\nWhat's the next step?"
            })

        # Call LLM via LiteLLM
        try:
            from litellm import completion

            self.logger.info(f"Calling {self.model} via LiteLLM (attempt {self.code_attempts})...")

            # Build kwargs for LiteLLM
            kwargs = {
                "model": self.model,
                "messages": self.conversation_history,
                "temperature": self.temperature,
                "max_tokens": self.max_tokens,
            }

            # Add optional API configuration
            if self.api_key:
                kwargs["api_key"] = self.api_key
            if self.api_base:
                kwargs["api_base"] = self.api_base

            response = completion(**kwargs)

            # Extract generated code
            code = response.choices[0].message.content.strip()

            # Add assistant response to history
            self.conversation_history.append({
                "role": "assistant",
                "content": code
            })

            # Clean up code (remove markdown if LLM added it despite instructions)
            code = self._extract_code_from_response(code)

            self.logger.debug(f"Generated code ({len(code)} chars):\n{code[:200]}...")

            return code

        except Exception as e:
            self.logger.error(f"Error calling LLM: {e}")
            import traceback
            self.logger.error(traceback.format_exc())
            # Return a safe fallback
            return f'print("Error generating code: {e}")'

    def _extract_code_from_response(self, response: str) -> str:
        """
        Extract Python code from LLM response.

        Handles cases where LLM wraps code in markdown blocks despite instructions.

        Args:
            response: Raw LLM response

        Returns:
            Cleaned Python code
        """
        import re

        # Check if wrapped in markdown code block
        code_block_pattern = r"```(?:python)?\s*\n(.*?)\n```"
        matches = re.findall(code_block_pattern, response, re.DOTALL)

        if matches:
            # Extract code from first code block
            return matches[0].strip()

        # No markdown wrapper, return as-is
        return response.strip()

    def _should_finish(self, state: CodeState) -> bool:
        """
        Decide if task is complete based on state.

        Uses Azure OpenAI to determine if the task has been accomplished.

        Args:
            state: Current CodeState

        Returns:
            True if task appears complete
        """
        # If execution failed repeatedly, might want to give up
        if not state.execution_success and self.code_attempts >= 5:
            self.logger.warning("Multiple failures, considering task failed")
            return True

        # Ask LLM if task is complete (using same model)
        try:
            from litellm import completion

            check_prompt = f"""Given this task:
{self.instruction}

Current PowerPoint state:
{self._format_state(state)}

Has the task been fully completed? Answer with just "YES" or "NO" followed by a brief reason."""

            kwargs = {
                "model": self.model,
                "messages": [
                    {"role": "system", "content": "You are evaluating if a PowerPoint manipulation task has been completed."},
                    {"role": "user", "content": check_prompt}
                ],
                "temperature": 0.3,
                "max_tokens": 100,
            }

            if self.api_key:
                kwargs["api_key"] = self.api_key
            if self.api_base:
                kwargs["api_base"] = self.api_base

            response = completion(**kwargs)

            answer = response.choices[0].message.content.strip().upper()

            if answer.startswith("YES"):
                self.logger.info(f"LLM says task is complete: {answer}")
                return True

            return False

        except Exception as e:
            self.logger.error(f"Error checking completion: {e}")
            # Fallback to simple heuristic
            return state.execution_success and "completed" in state.output.lower()

    def reset(self) -> None:
        """Reset agent state for new task."""
        super().reset()
        self.conversation_history = []
        self.code_attempts = 0


# Example: Simple rule-based agent
class RuleBasedCodeAgent(CodeAgent):
    """
    Simple rule-based agent for testing.

    Follows predefined rules to generate code.
    """

    def __init__(self, config: dict[str, Any] | None = None):
        super().__init__(config)
        self.step_count = 0

    def _generate_code(self, state_info: str) -> str:
        """Generate code using simple rules."""
        self.step_count += 1

        # Example: Add text to first slide
        if self.step_count == 1:
            return """
# Add title to first slide
if len(prs.slides) > 0:
    slide = prs.slides[0]
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Updated Title"
print("Added title")
"""

        # Example: Add new slide
        elif self.step_count == 2:
            return """
# Add new slide
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "New Slide"
print("Added new slide")
"""

        else:
            # Done
            return 'print("Task completed")'

    def _should_finish(self, state: CodeState) -> bool:
        """Finish after 3 steps."""
        return self.step_count >= 3
