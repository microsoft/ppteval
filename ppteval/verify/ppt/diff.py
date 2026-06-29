import hashlib
import os
import re
import xml.etree.ElementTree as ET
import xmltodict
import zipfile
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from ppteval.verify.ppt.slide_screenshots_utils import SlideScreenshotGenerator

# Note: If making changes to the dataclasses in this file,
# you may need to update the prompts accordingly.


@dataclass
class AnimationEffect:
    """Represents a single animation effect"""

    slide_id: str
    element_id: str
    effect_type: str
    trigger: str
    delay: float
    duration: float
    order: int
    # Note, element_text may sometimes be a superset of the finer grained text actually animated.
    element_text: Optional[str] = None
    element_type: Optional[str] = None

    def to_dict(self) -> Dict[str, Any]:
        return {
            "slide_id": self.slide_id,
            "element_id": self.element_id,
            "effect_type": self.effect_type,
            "trigger": self.trigger,
            "delay": self.delay,
            "duration": self.duration,
            "order": self.order,
            "element_text": self.element_text,
            "element_type": self.element_type,
        }


@dataclass
class SlideTransition:
    """Represents a slide transition"""

    slide_id: str
    transition_type: str
    duration: float
    direction: Optional[str] = None

    def to_dict(self) -> Dict[str, Any]:
        return {
            "slide_id": self.slide_id,
            "transition_type": self.transition_type,
            "duration": self.duration,
            "direction": self.direction,
        }


@dataclass
class Slide:
    """Represents a slide with its metadata"""

    slide_id: str
    slide_number: int
    title: Optional[str] = None
    layout_type: Optional[str] = None
    element_count: int = 0
    notes: Optional[str] = None
    content_hash: Optional[str] = None
    _serialized_elements: Optional[str] = None  # Internal use for deep comparison

    def to_dict(self) -> Dict[str, Any]:
        return {
            "slide_id": self.slide_id,
            "slide_number": self.slide_number,
            "title": self.title,
            "layout_type": self.layout_type,
            "element_count": self.element_count,
            "notes": self.notes,
            "content_hash": self.content_hash
        }

    def __repr__(self):
        return f"Slide(slide_number={self.slide_number}, title={self.title}, layout_type={self.layout_type}, element_count={self.element_count}, notes={self.notes}, content_hash={self.content_hash})"


@dataclass
class PowerPointDiff:
    """Container for PowerPoint differences"""

    added_animations: List[AnimationEffect]
    removed_animations: List[AnimationEffect]
    modified_animations: List[Tuple[AnimationEffect, AnimationEffect]]
    added_transitions: List[SlideTransition]
    removed_transitions: List[SlideTransition]
    modified_transitions: List[Tuple[SlideTransition, SlideTransition]]
    added_slides: List[Slide]
    removed_slides: List[Slide]
    modified_slides: List[Tuple[Slide, Slide]]

    def to_dict(self) -> Dict[str, Any]:
        return {
            "added_animations": [anim.to_dict() for anim in self.added_animations],
            "removed_animations": [anim.to_dict() for anim in self.removed_animations],
            "modified_animations": [[old.to_dict(), new.to_dict()] for old, new in self.modified_animations],
            "added_transitions": [trans.to_dict() for trans in self.added_transitions],
            "removed_transitions": [trans.to_dict() for trans in self.removed_transitions],
            "modified_transitions": [[old.to_dict(), new.to_dict()] for old, new in self.modified_transitions],
            "added_slides": [slide.to_dict() for slide in self.added_slides],
            "removed_slides": [slide.to_dict() for slide in self.removed_slides],
            "modified_slides": [[old.to_dict(), new.to_dict()] for old, new in self.modified_slides],
        }


@dataclass
class SlideScreenshot:
    """Represents a slide screenshot"""

    slide_number: int
    image_path: str
    slide_id: Optional[str] = None


class PowerPointDiffEvaluator:
    """Evaluates differences between PowerPoint files focusing on animations and transitions"""

    def __init__(self):
        self.namespaces = {
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        }
        self.slide_elements_cache = {}  # Cache for slide elements
        self.screenshot_generator = SlideScreenshotGenerator()

    def extract_animations_and_transitions(self, pptx_path: str) -> Tuple[List[AnimationEffect], List[SlideTransition]]:
        """Extract animations and transitions from a PowerPoint file"""
        animations = []
        transitions = []

        try:
            prs = Presentation(pptx_path)
            with zipfile.ZipFile(pptx_path, "r") as pptx_zip:
                for i, slide in enumerate(prs.slides):
                    try:
                        # Use python-pptx slide_id (same as in extract_slides)
                        slide_id = str(slide.slide_id)

                        # Get slide XML path
                        partname_obj = slide.part.partname  # type: ignore[attr-defined]
                        slide_path = str(partname_obj).lstrip("/")

                        # Read slide XML
                        slide_xml = pptx_zip.read(slide_path)
                        slide_root = ET.fromstring(slide_xml)

                        # Cache slide elements for text lookup
                        self.slide_elements_cache[slide_id] = self._build_element_map(slide_root)

                        # Extract transitions
                        transition = self._extract_transition(slide_root, slide_id)
                        if transition:
                            transitions.append(transition)

                        # Extract animations
                        slide_animations = self._extract_animations(slide_root, slide_id)
                        animations.extend(slide_animations)

                    except Exception as e:
                        print(f"Warning: Error processing slide animations/transitions: {e}")
                        continue

        except Exception as e:
            print(f"Error processing {pptx_path}: {e}")

        return animations, transitions

    def extract_slides(self, pptx_path: str) -> List[Slide]:
        """Extract slide metadata from a PowerPoint file"""
        slides = []

        try:
            prs = Presentation(pptx_path)
            with zipfile.ZipFile(pptx_path, "r") as pptx_zip:
                for i, slide in enumerate(prs.slides):
                    try:
                        # partname like '/ppt/slides/slide1.xml' (used only for logging).
                        # NOTE: python-pptx renumbers partnames virtually on load
                        # based on slide position in sldIdLst, so the partname
                        # may not match the on-disk filename in the zip. We must
                        # read the slide's XML directly from slide.part.blob to
                        # get the correct content for THIS slide.
                        partname_obj = slide.part.partname  # type: ignore[attr-defined]
                        partname_str = str(partname_obj)
                        slide_path = partname_str.lstrip("/")

                        # Use python-pptx slide_id as requested
                        slide_id = str(slide.slide_id)

                        # Slide number is the 1-based index
                        slide_number = i + 1

                        # Read slide XML directly from the loaded part (correct
                        # content even when python-pptx has renumbered partnames).
                        slide_xml = slide.part.blob  # type: ignore[attr-defined]
                        slide_root = ET.fromstring(slide_xml)

                        # Compute content hash of the raw slide XML to detect any change
                        content_hash = hashlib.sha256(slide_xml).hexdigest()

                        # Resolve the actual referenced slideLayout deterministically via python-pptx.
                        # This avoids the fuzzy placeholder-signature matcher in _extract_layout_type,
                        # whose result can flip across resaves of an unchanged slide.
                        layout_name_ref: Optional[str] = None
                        try:
                            layout_name_ref = slide.slide_layout.name  # type: ignore[attr-defined]
                        except Exception:
                            layout_name_ref = None

                        # Extract notes text via slide relationship to notesSlide
                        notes_text: Optional[str] = None
                        try:
                            for rel in slide.part.rels.values():  # type: ignore[attr-defined]
                                if rel.reltype == RT.NOTES_SLIDE:
                                    notes_xml = rel.target_part.blob
                                    notes_root = ET.fromstring(notes_xml)
                                    text_elements = notes_root.findall(".//a:t", self.namespaces)
                                    collected: List[str] = []
                                    for t in text_elements:
                                        if t.text:
                                            collected.append(t.text)
                                    joined = " ".join(collected).strip()
                                    notes_text = joined if joined else None
                                    break
                        except Exception as e1:
                            notes_text = None
                            print("Hit an exception: ", e1)

                        slide_meta = self._extract_slide_metadata(slide_root, slide_id, slide_number, notes_text, pptx_zip, layout_name_ref)
                        slide_meta.content_hash = content_hash
                        slide_meta._serialized_elements = self.slide_xml_to_json(slide_xml)
                        slides.append(slide_meta)
                    except Exception as e:
                        print(f"Warning: Error processing slide via python-pptx: {e}")
                        continue
            return slides

        except Exception as e:
            print(f"Error processing slides from {pptx_path}: {e}")

        return slides

    def _extract_slide_metadata(
        self,
        slide_root: ET.Element,
        slide_id: str,
        slide_number: int,
        notes_text: Optional[str] = None,
        pptx_zip: Optional[zipfile.ZipFile] = None,
        layout_name_ref: Optional[str] = None,
    ) -> Slide:
        """Extract metadata from a single slide"""
        # Extract title
        title = self._extract_slide_title(slide_root)

        # Extract layout type: prefer the deterministically-resolved layout name (from the
        # slide's actual slideLayout relationship). Fall back to the heuristic matcher only
        # if the deterministic lookup wasn't available.
        if layout_name_ref:
            layout_type: Optional[str] = layout_name_ref
        else:
            layout_type = self._extract_layout_type(slide_root, pptx_zip if pptx_zip else None)

        # Count elements
        element_count = self._count_slide_elements(slide_root)

        # Extract notes (if available): use relationship-derived notes only
        notes = notes_text

        return Slide(
            slide_id=slide_id,
            slide_number=slide_number,
            title=title,
            layout_type=layout_type,
            element_count=element_count,
            notes=notes,
        )

    def _extract_slide_title(self, slide_root: ET.Element) -> Optional[str]:
        """Extract title from slide"""
        # Look for title placeholder
        title_elements = slide_root.findall(".//p:sp", self.namespaces)

        for shape in title_elements:
            # Check if this is a title placeholder
            nv_sp_pr = shape.find(".//p:nvSpPr", self.namespaces)
            if nv_sp_pr is not None:
                ph = nv_sp_pr.find(".//p:ph", self.namespaces)
                if ph is not None:
                    ph_type = ph.get("type")
                    if ph_type in ["title", "ctrTitle"]:
                        # Extract text from this shape
                        return self._extract_text_from_shape(shape)

                # Also check by name
                c_nv_pr = nv_sp_pr.find(".//p:cNvPr", self.namespaces)
                if c_nv_pr is not None:
                    name = c_nv_pr.get("name", "").lower()
                    if "title" in name:
                        return self._extract_text_from_shape(shape)

        return None

    def _extract_layout_type(self, slide_root: ET.Element, pptx: Optional[zipfile.ZipFile] = None) -> Optional[str]:
        """Extract layout type from slide by reading the actual layout reference"""

        try:
            # Look for all layout files and match by placeholders
            if pptx is None:
                layout_files = []
            else:
                layout_files = [f for f in pptx.namelist() if f.startswith("ppt/slideLayouts/slideLayout") and f.endswith(".xml")]

            if layout_files:
                # Get placeholders from current slide
                slide_placeholders = slide_root.findall(".//p:ph", self.namespaces)
                slide_ph_signature = []

                for ph in slide_placeholders:
                    ph_type = ph.get("type", "none")
                    ph_idx = ph.get("idx", "0")
                    slide_ph_signature.append((ph_type, ph_idx))

                # Sort for consistent comparison
                slide_ph_signature.sort()

                best_match = None
                best_match_score = 0

                # Try to match with layout files
                for layout_file in layout_files:
                    try:
                        layout_xml = pptx.read(layout_file)
                        layout_root = ET.fromstring(layout_xml)

                        # Extract layout name - look for name attribute in cSld
                        cSld = layout_root.find('.//p:cSld', self.namespaces)
                        layout_name = None
                        if cSld is not None:
                            layout_name = cSld.get('name')

                        # If no name in cSld, try to find it in other elements
                        if not layout_name:
                            # Sometimes the name is in other locations
                            for elem in layout_root.iter():
                                if elem.get('name'):
                                    layout_name = elem.get('name')
                                    break

                        if not layout_name:
                            # Extract from filename as fallback
                            layout_name = layout_file.split('/')[-1].replace('.xml', '').replace('slideLayout', 'Layout ')

                        # Get layout placeholders for matching
                        layout_placeholders = layout_root.findall(".//p:ph", self.namespaces)
                        layout_ph_signature = []

                        for ph in layout_placeholders:
                            ph_type = ph.get("type", "none")
                            ph_idx = ph.get("idx", "0")
                            layout_ph_signature.append((ph_type, ph_idx))

                        layout_ph_signature.sort()

                        # Calculate match score with more flexible matching
                        match_score = 0.0

                        if len(slide_ph_signature) == 0 and len(layout_ph_signature) == 0:
                            match_score = 1.0  # Both empty
                        elif len(slide_ph_signature) == 0 or len(layout_ph_signature) == 0:
                            match_score = 0.0
                        else:
                            # Method 1: Exact signature match (highest priority)
                            if slide_ph_signature == layout_ph_signature:
                                match_score = 1.0
                            else:
                                # Method 2: Count matching placeholders (flexible)
                                slide_set = set(slide_ph_signature)
                                layout_set = set(layout_ph_signature)

                                # Count matches
                                matches = len(slide_set.intersection(layout_set))

                                # Use the larger set as denominator for more flexible matching
                                total = max(len(slide_set), len(layout_set))
                                if total > 0:
                                    match_score = matches / total

                                # Method 3: Check if slide placeholders are subset of layout
                                if len(slide_set) <= len(layout_set):
                                    subset_matches = len(slide_set.intersection(layout_set))
                                    if subset_matches == len(slide_set):
                                        # All slide placeholders are in layout - good match
                                        match_score = max(match_score, 0.9)

                                # Method 4: Special handling for common patterns
                                slide_types = {ph[0] for ph in slide_ph_signature}
                                layout_types = {ph[0] for ph in layout_ph_signature}

                                # Check for title + content patterns
                                if "title" in slide_types and "title" in layout_types:
                                    if slide_types.intersection(layout_types):
                                        match_score = max(match_score, 0.7)

                        # Update best match if this is better
                        if match_score > best_match_score:
                            best_match_score = match_score
                            best_match = layout_name

                    except Exception:
                        continue

                # If we found a good match, return it (lowered threshold for better matching)
                if best_match and best_match_score >= 0.6:  # Reduced from 0.8 to 0.6
                    return best_match

            # Fallback to the placeholder-based detection if layout matching fails
            placeholders = slide_root.findall(".//p:ph", self.namespaces)

            ph_types = set()
            ph_indices = []
            for ph in placeholders:
                ph_type = ph.get("type")
                ph_idx = ph.get("idx")
                if ph_type:
                    ph_types.add(ph_type)
                if ph_idx:
                    ph_indices.append(int(ph_idx))

            # Count content placeholders more accurately
            content_placeholders = 0
            title_placeholders = 0

            for ph in placeholders:
                ph_type = ph.get("type")
                ph_idx = ph.get("idx")

                if ph_type in ["title", "ctrTitle"]:
                    title_placeholders += 1
                elif ph_type in ["body", "obj", "pic", "chart", "tbl", "media"]:
                    content_placeholders += 1
                elif ph_type is None and ph_idx:
                    # Some content placeholders might not have explicit types
                    # Check if it's likely a content placeholder based on index
                    idx = int(ph_idx)
                    if idx > 0:  # Title is usually idx=0, content starts from 1
                        content_placeholders += 1

            # Determine layout based on placeholder analysis
            if title_placeholders > 0 and content_placeholders >= 3:
                return "Three Content"
            elif title_placeholders > 0 and content_placeholders == 2:
                return "Two Content"
            elif "title" in ph_types and "body" in ph_types and content_placeholders == 1:
                return "Title and Content"
            elif "title" in ph_types and not ph_types - {"title"}:
                return "Title Only"
            elif "ctrTitle" in ph_types and "subTitle" in ph_types:
                return "Title Slide"
            elif content_placeholders > 1:
                return f"{content_placeholders} Content"
            elif not ph_types:
                return "Blank"
            else:
                return "Custom Layout"

        except Exception as e:
            print(f"Error extracting layout type: {e}")
            return "Unknown Layout"

    def _count_slide_elements(self, slide_root: ET.Element) -> int:
        """Count the number of elements on a slide"""
        # Count shapes
        shapes = slide_root.findall(".//p:sp", self.namespaces)
        return len(shapes)

    # Removed: _extract_slide_notes (no fallback; notes resolved via python-pptx relationships)

    def _build_element_map(self, slide_root: ET.Element) -> Dict[str, Dict[str, str]]:
        """Build a map of element IDs to their text content and types"""
        element_map = {}

        # Find all shapes in the slide
        shapes = slide_root.findall(".//p:sp", self.namespaces)

        for shape in shapes:
            # Get shape ID
            nv_sp_pr = shape.find(".//p:nvSpPr", self.namespaces)
            if nv_sp_pr is not None:
                c_nv_pr = nv_sp_pr.find(".//p:cNvPr", self.namespaces)
                if c_nv_pr is not None:
                    shape_id = c_nv_pr.get("id")
                    shape_name = c_nv_pr.get("name", "")

                    if shape_id:
                        # Extract text content
                        text_content = self._extract_text_from_shape(shape)

                        # Determine element type
                        element_type = self._determine_element_type(shape, shape_name)

                        element_map[shape_id] = {
                            "text": text_content,
                            "type": element_type,
                            "name": shape_name,
                        }

        # Also check for group shapes
        groups = slide_root.findall(".//p:grpSp", self.namespaces)
        for group in groups:
            group_shapes = group.findall(".//p:sp", self.namespaces)
            for shape in group_shapes:
                nv_sp_pr = shape.find(".//p:nvSpPr", self.namespaces)
                if nv_sp_pr is not None:
                    c_nv_pr = nv_sp_pr.find(".//p:cNvPr", self.namespaces)
                    if c_nv_pr is not None:
                        shape_id = c_nv_pr.get("id")
                        shape_name = c_nv_pr.get("name", "")

                        if shape_id:
                            text_content = self._extract_text_from_shape(shape)
                            element_type = self._determine_element_type(shape, shape_name)

                            element_map[shape_id] = {
                                "text": text_content,
                                "type": element_type,
                                "name": shape_name,
                            }

        return element_map

    def _extract_text_from_shape(self, shape: ET.Element) -> Optional[str]:
        """Extract all text content from a shape"""
        text_parts = []

        # Find all text elements
        text_elements = shape.findall(".//a:t", self.namespaces)
        for text_elem in text_elements:
            if text_elem.text:
                text_parts.append(text_elem.text)

        # Join all text parts
        full_text = "".join(text_parts).strip()

        # Clean up excessive whitespace but preserve line breaks
        full_text = re.sub(r"[ \t]+", " ", full_text)  # Only collapse spaces/tabs, preserve newlines

        return full_text if full_text else None

    def _get_element_text(self, slide_root: ET.Element, spid: str) -> Optional[str]:
        """Get text content for a specific shape ID"""
        try:
            # Find the shape with the matching spid
            shapes = slide_root.findall(".//p:sp", self.namespaces)
            for shape in shapes:
                nvSpPr = shape.find(".//p:cNvPr", self.namespaces)
                if nvSpPr is not None and nvSpPr.get("id") == spid:
                    return self._extract_text_from_shape(shape)
        except Exception:
            pass
        return None

    def _determine_element_type(self, shape: ET.Element, shape_name: str) -> str:
        """Determine the type of element based on shape properties"""
        # Check if it's a text box
        if "TextBox" in shape_name or "Text Box" in shape_name:
            return "textbox"

        # Check if it's a title
        if "Title" in shape_name:
            return "title"

        # Check if it's a content placeholder
        if "Content Placeholder" in shape_name or "Content" in shape_name:
            return "content"

        # Check if it has text body
        tx_body = shape.find(".//p:txBody", self.namespaces)
        if tx_body is not None:
            return "text_shape"

        # Check for specific shape types
        sp_pr = shape.find(".//p:spPr", self.namespaces)
        if sp_pr is not None:
            # Check for preset geometry
            prst_geom = sp_pr.find(".//a:prstGeom", self.namespaces)
            if prst_geom is not None:
                prst = prst_geom.get("prst", "")
                if prst == "rect":
                    return "rectangle"
                elif prst == "ellipse":
                    return "ellipse"
                elif prst in ["star5", "star6", "star8"]:
                    return "star"
                elif prst in ["rightArrow", "leftArrow", "upArrow", "downArrow"]:
                    return "arrow"
                else:
                    return f"shape_{prst}"

        return "shape"

    def _extract_transition(self, slide_root: ET.Element, slide_id: str) -> Optional[SlideTransition]:
        """Extract transition information from slide XML"""
        transition_elem = slide_root.find(".//p:transition", self.namespaces)
        if transition_elem is None:
            return None

        # Get transition type
        transition_type = "none"
        for child in transition_elem:
            if "}" in child.tag:
                # Extract the tag name without the namespace, e.g., 'fade' from '{...}fade'
                tag_name = child.tag.split("}")[-1]

                # Child elements of <p:transition> can be the transition type (e.g., <p:fade>)
                # or other properties (e.g., <p:soundAc> for sound).
                # We assume that any child that is not a known property is the transition type.
                if tag_name != "soundAc":
                    transition_type = tag_name
                    break  # Assume the first such element is the transition type

            # if child.tag.endswith('}fade'):
            #     transition_type = "fade"
            # elif child.tag.endswith('}push'):
            #     transition_type = "push"
            # elif child.tag.endswith('}wipe'):
            #     transition_type = "wipe"
            # elif child.tag.endswith('}cut'):
            #     transition_type = "cut"
            # elif child.tag.endswith('}dissolve'):
            #     transition_type = "dissolve"
            # elif child.tag.endswith('}morph'):
            #     transition_type = "morph"
            # # Add more transition types as needed

        # Get duration (in milliseconds, convert to seconds)
        duration_attr = transition_elem.get("dur", "500")
        duration = float(duration_attr) / 1000.0

        # Get direction if available
        direction = None
        for child in transition_elem:
            direction_attr = child.get("dir")
            if direction_attr:
                direction = direction_attr
                break

        return SlideTransition(
            slide_id=slide_id,
            transition_type=transition_type,
            duration=duration,
            direction=direction,
        )

    def _extract_animations(self, slide_root: ET.Element, slide_id: str) -> List[AnimationEffect]:
        """Extract animation effects from slide XML"""
        animations = []

        # Find timing root
        timing_root = slide_root.find(".//p:timing", self.namespaces)
        if timing_root is None:
            return animations

        parent_map = {c: p for p in timing_root.iter() for c in p}

        # Look for various animation elements
        # Enhanced detection: Look for specific animation presets
        animations.extend(self._detect_entrance_animations(timing_root, slide_root, slide_id))
        animations.extend(self._detect_exit_animations(timing_root, slide_root, slide_id))
        animations.extend(self._detect_emphasis_animations(timing_root, slide_root, slide_id))
        animations.extend(self._detect_motion_path_animations(timing_root, slide_root, slide_id))

        # Look for various other animation elements (fallback to original logic)
        animation_elements = []

        # Find animation effects
        animation_elements.extend(timing_root.findall(".//p:animEffect", self.namespaces))

        # Find animation motions
        animation_elements.extend(timing_root.findall(".//p:animMotion", self.namespaces))

        # Find animation colors
        animation_elements.extend(timing_root.findall(".//p:animClr", self.namespaces))

        # Find animation rotations
        animation_elements.extend(timing_root.findall(".//p:animRot", self.namespaces))

        # Find animation scales
        animation_elements.extend(timing_root.findall(".//p:animScale", self.namespaces))

        # Find set animations (appear/disappear)
        animation_elements.extend(timing_root.findall(".//p:set", self.namespaces))

        # Find animation groups/sequences
        par_elements = timing_root.findall(".//p:par", self.namespaces)
        for par in par_elements:
            # Look for animations within parallel groups
            animation_elements.extend(par.findall(".//p:animEffect", self.namespaces))
            animation_elements.extend(par.findall(".//p:animMotion", self.namespaces))
            animation_elements.extend(par.findall(".//p:animClr", self.namespaces))
            animation_elements.extend(par.findall(".//p:animRot", self.namespaces))
            animation_elements.extend(par.findall(".//p:animScale", self.namespaces))
            animation_elements.extend(par.findall(".//p:set", self.namespaces))

        # Parse each animation element (only if not already detected as specific preset)
        existing_spids = {anim.element_id for anim in animations}

        for i, anim_elem in enumerate(animation_elements):
            animation = self._parse_animation_effect(anim_elem, slide_id, i, parent_map)
            if animation and animation.element_id not in existing_spids:
                animations.append(animation)

        return animations

    def _detect_entrance_animations(self, timing_root: ET.Element, slide_root: ET.Element, slide_id: str) -> List[AnimationEffect]:
        """Detect entrance animations (appear, fly in, fade in, etc.)"""
        animations = []

        # Track animations we've already added to avoid duplicates
        added_animations = set()  # Will store (slide_id, element_id, effect_type) tuples

        # Common entrance animation presets
        entrance_presets = {
            # Fly In animations
            ('2', 'entr', '8'): 'fly_in_from_left',
            ('2', 'entr', '2'): 'fly_in_from_right',
            ('2', 'entr', '4'): 'fly_in_from_top',
            ('2', 'entr', '6'): 'fly_in_from_bottom',
            ('2', 'entr', '1'): 'fly_in_from_top_left',
            ('2', 'entr', '3'): 'fly_in_from_top_right',
            ('2', 'entr', '7'): 'fly_in_from_bottom_left',
            ('2', 'entr', '9'): 'fly_in_from_bottom_right',

            # Fade animations
            ('1', 'entr', '0'): 'fade_in',

            # Appear animation
            ('0', 'entr', '0'): 'appear',

            # Wipe animations
            ('5', 'entr', '8'): 'wipe_from_left',
            ('5', 'entr', '2'): 'wipe_from_right',
            ('5', 'entr', '4'): 'wipe_from_top',
            ('5', 'entr', '6'): 'wipe_from_bottom',

            # Split animations
            ('13', 'entr', '10'): 'split_horizontal_in',
            ('13', 'entr', '5'): 'split_vertical_in',

            # Blinds animations
            ('3', 'entr', '5'): 'blinds_horizontal',
            ('3', 'entr', '10'): 'blinds_vertical',

            # Box animations
            ('4', 'entr', '5'): 'box_in',
            ('4', 'entr', '0'): 'box_out',

            # Checkerboard animations
            ('6', 'entr', '5'): 'checkerboard_across',
            ('6', 'entr', '10'): 'checkerboard_down',

            # Dissolve animation
            ('8', 'entr', '0'): 'dissolve_in',

            # Random Bars
            ('9', 'entr', '5'): 'random_bars_horizontal',
            ('9', 'entr', '10'): 'random_bars_vertical',

            # Strips animations
            ('10', 'entr', '1'): 'strips_left_up',
            ('10', 'entr', '3'): 'strips_right_up',
            ('10', 'entr', '7'): 'strips_left_down',
            ('10', 'entr', '9'): 'strips_right_down',

            # Zoom animations
            ('16', 'entr', '0'): 'zoom_in',
            ('16', 'entr', '1'): 'zoom_in_from_screen_center',
            ('16', 'entr', '2'): 'zoom_in_slightly',

            # Stretch animations
            ('14', 'entr', '8'): 'stretch_from_left',
            ('14', 'entr', '2'): 'stretch_from_right',
            ('14', 'entr', '4'): 'stretch_from_top',
            ('14', 'entr', '6'): 'stretch_from_bottom',

            # Swivel animation
            ('15', 'entr', '0'): 'swivel',

            # Spiral animation
            ('11', 'entr', '0'): 'spiral_in',
        }

        for (preset_id, preset_class, preset_subtype), effect_name in entrance_presets.items():
            ctn_elements = timing_root.findall(f'.//p:cTn[@presetID="{preset_id}"][@presetClass="{preset_class}"][@presetSubtype="{preset_subtype}"]', self.namespaces)

            for ctn in ctn_elements:
                target_elements = ctn.findall('.//p:spTgt', self.namespaces)

                for target in target_elements:
                    spid = target.get('spid')

                    # Create unique key for this animation to avoid duplicates
                    anim_key = (slide_id, spid, effect_name)
                    if anim_key in added_animations:
                        continue  # Skip duplicate

                    text_elements = target.findall('.//p:txEl', self.namespaces)
                    is_text_animation = len(text_elements) > 0

                    # Verify this is actually the correct animation by checking motion patterns for fly-in animations
                    if effect_name.startswith('fly_in_from_'):
                        if not self._verify_fly_in_direction(ctn, effect_name):
                            continue

                    element_text = self._get_element_text(slide_root, spid)

                    animation = AnimationEffect(
                        slide_id=slide_id,
                        element_id=spid,
                        effect_type=effect_name,
                        trigger='click',  # Default trigger for entrance animations
                        delay=0.0,  # Default delay
                        duration=0.5,  # Default duration
                        order=len(animations),
                        element_text=element_text,
                        element_type='text' if is_text_animation else 'shape'
                    )
                    animations.append(animation)
                    added_animations.add(anim_key)  # Track this animation as added

        return animations

    def _detect_exit_animations(self, timing_root: ET.Element, slide_root: ET.Element, slide_id: str) -> List[AnimationEffect]:
        """Detect exit animations (disappear, fly out, fade out, etc.)"""
        animations = []

        # Track animations we've already added to avoid duplicates
        added_animations = set()  # Will store (slide_id, element_id, effect_type) tuples

        # Common exit animation presets
        exit_presets = {
            # Fly Out animations
            ('2', 'exit', '2'): 'fly_out_to_right',
            ('2', 'exit', '8'): 'fly_out_to_left',
            ('2', 'exit', '4'): 'fly_out_to_top',
            ('2', 'exit', '6'): 'fly_out_to_bottom',

            # Fade Out
            ('1', 'exit', '0'): 'fade_out',

            # Disappear
            ('0', 'exit', '0'): 'disappear',

            # Wipe Out animations
            ('5', 'exit', '2'): 'wipe_out_to_right',
            ('5', 'exit', '8'): 'wipe_out_to_left',
            ('5', 'exit', '4'): 'wipe_out_to_top',
            ('5', 'exit', '6'): 'wipe_out_to_bottom',

            # Split Out
            ('13', 'exit', '10'): 'split_horizontal_out',
            ('13', 'exit', '5'): 'split_vertical_out',

            # Box Out
            ('4', 'exit', '0'): 'box_out',

            # Zoom Out
            ('16', 'exit', '0'): 'zoom_out',

            # Spiral Out
            ('11', 'exit', '0'): 'spiral_out',
        }

        for (preset_id, preset_class, preset_subtype), effect_name in exit_presets.items():
            ctn_elements = timing_root.findall(f'.//p:cTn[@presetID="{preset_id}"][@presetClass="{preset_class}"][@presetSubtype="{preset_subtype}"]', self.namespaces)

            for ctn in ctn_elements:
                target_elements = ctn.findall('.//p:spTgt', self.namespaces)

                for target in target_elements:
                    spid = target.get('spid')

                    # Create unique key for this animation to avoid duplicates
                    anim_key = (slide_id, spid, effect_name)
                    if anim_key in added_animations:
                        continue  # Skip duplicate

                    text_elements = target.findall('.//p:txEl', self.namespaces)
                    is_text_animation = len(text_elements) > 0

                    element_text = self._get_element_text(slide_root, spid)

                    animation = AnimationEffect(
                        slide_id=slide_id,
                        element_id=spid,
                        effect_type=effect_name,
                        trigger='on_click',  # Default trigger for exit animations
                        delay=0.0,  # Default delay
                        duration=0.5,  # Default duration
                        order=len(animations),
                        element_text=element_text,
                        element_type='text' if is_text_animation else 'shape'
                    )
                    animations.append(animation)
                    added_animations.add(anim_key)  # Track this animation as added

        return animations

    def _detect_emphasis_animations(self, timing_root: ET.Element, slide_root: ET.Element, slide_id: str) -> List[AnimationEffect]:
        """Detect emphasis animations (pulse, color emphasis, grow/shrink, etc.)"""
        animations = []

        # Track animations we've already added to avoid duplicates
        added_animations = set()  # Will store (slide_id, element_id, effect_type) tuples

        # Common emphasis animation presets
        emphasis_presets = {
            # Color emphasis
            ('19', 'emph', '0'): 'color_emphasis',

            # Pulse/Grow Shrink
            ('1', 'emph', '0'): 'pulse',
            ('2', 'emph', '0'): 'grow_shrink',

            # Spin
            ('3', 'emph', '0'): 'spin',

            # Transparency
            ('18', 'emph', '0'): 'transparency',

            # Bold Flash
            ('4', 'emph', '0'): 'bold_flash',

            # Bold Reveal
            ('5', 'emph', '0'): 'bold_reveal',

            # Wave
            ('6', 'emph', '0'): 'wave',

            # Lighten
            ('7', 'emph', '0'): 'lighten',

            # Darken
            ('8', 'emph', '0'): 'darken',

            # Complementary Color
            ('9', 'emph', '0'): 'complementary_color',

            # Complementary Color 2
            ('10', 'emph', '0'): 'complementary_color_2',

            # Contrasting Color
            ('11', 'emph', '0'): 'contrasting_color',

            # Brush on Color
            ('12', 'emph', '0'): 'brush_on_color',

            # Reveal Underline
            ('13', 'emph', '0'): 'reveal_underline',

            # Underline
            ('14', 'emph', '0'): 'underline',

            # Style Emphasis
            ('15', 'emph', '0'): 'style_emphasis',

            # Font Color
            ('16', 'emph', '0'): 'font_color',

            # Teeter
            ('17', 'emph', '0'): 'teeter',

            # Flicker
            ('20', 'emph', '0'): 'flicker',
        }

        for (preset_id, preset_class, preset_subtype), effect_name in emphasis_presets.items():
            ctn_elements = timing_root.findall(f'.//p:cTn[@presetID="{preset_id}"][@presetClass="{preset_class}"][@presetSubtype="{preset_subtype}"]', self.namespaces)

            for ctn in ctn_elements:
                target_elements = ctn.findall('.//p:spTgt', self.namespaces)

                for target in target_elements:
                    spid = target.get('spid')

                    # Create unique key for this animation to avoid duplicates
                    anim_key = (slide_id, spid, effect_name)
                    if anim_key in added_animations:
                        continue  # Skip duplicate

                    text_elements = target.findall('.//p:txEl', self.namespaces)
                    is_text_animation = len(text_elements) > 0

                    element_text = self._get_element_text(slide_root, spid)

                    animation = AnimationEffect(
                        slide_id=slide_id,
                        element_id=spid,
                        effect_type=effect_name,
                        trigger='with_previous',  # Default trigger for emphasis animations
                        delay=0.0,  # Default delay
                        duration=1.0,  # Default duration
                        order=len(animations),
                        element_text=element_text,
                        element_type='text' if is_text_animation else 'shape'
                    )
                    animations.append(animation)
                    added_animations.add(anim_key)  # Track this animation as added

        return animations

    def _detect_motion_path_animations(self, timing_root: ET.Element, slide_root: ET.Element, slide_id: str) -> List[AnimationEffect]:
        """Detect motion path animations"""
        animations = []

        # Look for animMotion elements which indicate motion path animations
        motion_elements = timing_root.findall('.//p:animMotion', self.namespaces)

        for motion in motion_elements:
            # Get the behavior element to find target
            cBhvr = motion.find('.//p:cBhvr', self.namespaces)
            if cBhvr is None:
                continue

            target_elem = cBhvr.find('.//p:tgtEl//p:spTgt', self.namespaces)
            if target_elem is None:
                continue

            spid = target_elem.get('spid')
            text_elements = target_elem.findall('.//p:txEl', self.namespaces)
            is_text_animation = len(text_elements) > 0

            # Analyze the path to determine motion type
            path = motion.get('path', '')
            motion_type = self._analyze_motion_path(path)

            element_text = self._get_element_text(slide_root, spid)

            animation = AnimationEffect(
                slide_id=slide_id,
                element_id=spid,
                effect_type=motion_type,
                trigger='on_click',  # Default trigger for motion path animations
                delay=0.0,  # Default delay
                duration=2.0,  # Default duration for motion paths
                order=len(animations),
                element_text=element_text,
                element_type='text' if is_text_animation else 'shape'
            )
            animations.append(animation)

        return animations

    def _verify_fly_in_direction(self, ctn: ET.Element, expected_effect: str) -> bool:
        """Verify that a fly-in animation matches the expected direction"""
        anim_elements = ctn.findall('.//p:anim', self.namespaces)

        for anim in anim_elements:
            attr_names = anim.findall('.//p:attrName', self.namespaces)

            for attr in attr_names:
                if attr.text == 'ppt_x':  # X-position animation
                    tav_elements = anim.findall('.//p:tav', self.namespaces)

                    for tav in tav_elements:
                        if tav.get('tm') == '0':  # Starting position
                            val_elements = tav.findall('.//p:strVal', self.namespaces)
                            for val in val_elements:
                                start_val = val.get('val', '')

                                # Check if the starting position matches the expected direction
                                if expected_effect == 'fly_in_from_left' and '0-#ppt_w/2' in start_val:
                                    return True
                                elif expected_effect == 'fly_in_from_right' and '#ppt_w' in start_val:
                                    return True

                elif attr.text == 'ppt_y':  # Y-position animation
                    tav_elements = anim.findall('.//p:tav', self.namespaces)

                    for tav in tav_elements:
                        if tav.get('tm') == '0':  # Starting position
                            val_elements = tav.findall('.//p:strVal', self.namespaces)
                            for val in val_elements:
                                start_val = val.get('val', '')

                                # Check if the starting position matches the expected direction
                                if expected_effect == 'fly_in_from_top' and '0-#ppt_h/2' in start_val:
                                    return True
                                elif expected_effect == 'fly_in_from_bottom' and '#ppt_h' in start_val:
                                    return True

        return True  # Default to accepting if we can't verify (preset ID should be sufficient)

    def _analyze_motion_path(self, path: str) -> str:
        """Analyze a motion path string to determine the type of movement"""
        if not path:
            return 'motion_path_custom'

        path = path.lower()

        # Common motion path patterns
        if 'l' in path and 'c' in path:
            return 'motion_path_curved'
        elif 'l' in path:
            return 'motion_path_linear'
        elif 'c' in path:
            return 'motion_path_curve'
        elif 'a' in path:
            return 'motion_path_arc'
        else:
            return 'motion_path_custom'

    def _parse_animation_effect(
        self,
        effect_elem: ET.Element,
        slide_id: str,
        order: int,
        parent_map: Dict[ET.Element, ET.Element],
    ) -> Optional[AnimationEffect]:
        """Parse a single animation effect element"""
        try:
            # Determine animation type based on element tag
            tag_name = effect_elem.tag.split("}")[-1] if "}" in effect_elem.tag else effect_elem.tag

            # effect_type = "unknown"
            # if tag_name == "animEffect":
            #     effect_type = effect_elem.get('filter', 'effect')
            # elif tag_name == "animMotion":
            #     effect_type = "motion"
            # elif tag_name == "animClr":
            #     effect_type = "color"
            # elif tag_name == "animRot":
            #     effect_type = "rotation"
            # elif tag_name == "animScale":
            #     effect_type = "scale"
            # elif tag_name == "set":
            #     effect_type = "set"

            effect_type = tag_name

            # Get target element ID - try multiple approaches
            element_id = "unknown"

            # First try to find target element directly
            target_elem = effect_elem.find(".//p:tgtEl", self.namespaces)
            if target_elem is not None:
                # Try shape target
                sp_tgt = target_elem.find(".//p:spTgt", self.namespaces)
                if sp_tgt is not None:
                    element_id = sp_tgt.get("spid", "unknown")
                else:
                    # Try ink target
                    ink_tgt = target_elem.find(".//p:inkTgt", self.namespaces)
                    if ink_tgt is not None:
                        element_id = ink_tgt.get("spid", "unknown")
                    else:
                        # Try text target
                        txt_tgt = target_elem.find(".//p:txtTgt", self.namespaces)
                        if txt_tgt is not None:
                            element_id = txt_tgt.get("spid", "unknown")

            # Try to find target in parent container if not found directly
            if element_id == "unknown":
                parent = parent_map.get(effect_elem)
                while parent is not None and element_id == "unknown":
                    target_elem = parent.find(".//p:tgtEl", self.namespaces)
                    if target_elem is not None:
                        sp_tgt = target_elem.find(".//p:spTgt", self.namespaces)
                        if sp_tgt is not None:
                            element_id = sp_tgt.get("spid", "unknown")
                            break
                    parent = parent_map.get(parent)

            # Get timing information
            timing_elem = effect_elem.find(".//p:cTn", self.namespaces)
            delay = 0.0
            duration = 1.0
            if timing_elem is not None:
                delay_str = timing_elem.get("delay", "0")
                dur_str = timing_elem.get("dur", "1000")

                # Handle special duration values
                if dur_str == "indefinite":
                    duration = 0.0
                else:
                    try:
                        delay = float(delay_str) / 1000.0
                        duration = float(dur_str) / 1000.0
                    except ValueError:
                        delay = 0.0
                        duration = 1.0

            # Get trigger type - look in various places
            trigger = "onClick"  # Default

            # Look for condition elements
            cond_elem = effect_elem.find(".//p:cond", self.namespaces)
            if cond_elem is not None:
                trigger = cond_elem.get("evt", "onClick")
            else:
                # Look in parent elements
                parent = parent_map.get(effect_elem)
                while parent is not None and trigger == "onClick":
                    cond_elem = parent.find(".//p:cond", self.namespaces)
                    if cond_elem is not None:
                        trigger = cond_elem.get("evt", "onClick")
                        break
                    parent = parent_map.get(parent)

            # Map common trigger types
            trigger_map = {
                "onBegin": "withPrevious",
                "onEnd": "afterPrevious",
                "onClick": "onClick",
                "onDblClick": "onDoubleClick",
                "onMouseOver": "onMouseOver",
                "onMouseOut": "onMouseOut",
            }
            trigger = trigger_map.get(trigger, trigger)

            # Get text content and element type from cache
            element_text = None
            element_type = None
            if slide_id in self.slide_elements_cache and element_id in self.slide_elements_cache[slide_id]:
                element_info = self.slide_elements_cache[slide_id][element_id]
                element_text = element_info["text"]
                element_type = element_info["type"]

            return AnimationEffect(
                slide_id=slide_id,
                element_id=element_id,
                effect_type=effect_type,
                trigger=trigger,
                delay=delay,
                duration=duration,
                order=order,
                element_text=element_text,
                element_type=element_type,
            )

        except Exception as e:
            print(f"Error parsing animation effect: {e}")
            return None

    def compare_files(self, before_path: str, after_path: str) -> PowerPointDiff:
        """Compare two PowerPoint files and return differences"""
        # Extract animations and transitions from both files
        before_animations, before_transitions = self.extract_animations_and_transitions(before_path)
        after_animations, after_transitions = self.extract_animations_and_transitions(after_path)

        # Extract slides from both files
        before_slides = self.extract_slides(before_path)
        after_slides = self.extract_slides(after_path)

        # Compare animations
        added_animations, removed_animations, modified_animations = self._compare_animations(before_animations, after_animations)

        # Compare transitions
        added_transitions, removed_transitions, modified_transitions = self._compare_transitions(before_transitions, after_transitions)

        # Compare slides
        added_slides, removed_slides, modified_slides = self._compare_slides(before_slides, after_slides)

        return PowerPointDiff(
            added_animations=added_animations,
            removed_animations=removed_animations,
            modified_animations=modified_animations,
            added_transitions=added_transitions,
            removed_transitions=removed_transitions,
            modified_transitions=modified_transitions,
            added_slides=added_slides,
            removed_slides=removed_slides,
            modified_slides=modified_slides,
        )

    def _compare_animations(self, before: List[AnimationEffect], after: List[AnimationEffect]) -> Tuple[List[AnimationEffect], List[AnimationEffect], List[Tuple[AnimationEffect, AnimationEffect]]]:
        """Compare animation lists"""
        # Create lookup dictionaries
        before_dict = {self._animation_key(anim): anim for anim in before}
        after_dict = {self._animation_key(anim): anim for anim in after}

        # Find differences
        added = [anim for key, anim in after_dict.items() if key not in before_dict]
        removed = [anim for key, anim in before_dict.items() if key not in after_dict]

        # Find modified animations
        modified = []
        for key in before_dict:
            if key in after_dict:
                before_anim = before_dict[key]
                after_anim = after_dict[key]
                if not self._animations_equal(before_anim, after_anim):
                    modified.append((before_anim, after_anim))

        return added, removed, modified

    def _compare_transitions(self, before: List[SlideTransition], after: List[SlideTransition]) -> Tuple[List[SlideTransition], List[SlideTransition], List[Tuple[SlideTransition, SlideTransition]]]:
        """Compare transition lists"""
        # Create lookup dictionaries
        before_dict = {trans.slide_id: trans for trans in before}
        after_dict = {trans.slide_id: trans for trans in after}

        # Find differences
        added = [trans for slide_id, trans in after_dict.items() if slide_id not in before_dict]
        removed = [trans for slide_id, trans in before_dict.items() if slide_id not in after_dict]

        # Find modified transitions
        modified = []
        for slide_id in before_dict:
            if slide_id in after_dict:
                before_trans = before_dict[slide_id]
                after_trans = after_dict[slide_id]
                if not self._transitions_equal(before_trans, after_trans):
                    modified.append((before_trans, after_trans))

        return added, removed, modified

    def _compare_slides(self, before: List[Slide], after: List[Slide]) -> Tuple[List[Slide], List[Slide], List[Tuple[Slide, Slide]]]:
        """Compare slide metadata lists"""
        # Create lookup dictionaries
        before_dict = {slide.slide_id: slide for slide in before}
        after_dict = {slide.slide_id: slide for slide in after}

        # Find differences
        added = [slide for slide_id, slide in after_dict.items() if slide_id not in before_dict]
        removed = [slide for slide_id, slide in before_dict.items() if slide_id not in after_dict]

        # Find modified slides
        modified = []
        for slide_id in before_dict:
            if slide_id in after_dict:
                before_slide = before_dict[slide_id]
                after_slide = after_dict[slide_id]
                if not self._slides_equal(before_slide, after_slide):
                    modified.append((before_slide, after_slide))

        return added, removed, modified

    def _slides_equal(self, slide1: Slide, slide2: Slide) -> bool:
        """Check if two slides are equal"""
        # Prefer strong content hash comparison when available
        if slide1.content_hash and slide2.content_hash:
            if slide1.content_hash == slide2.content_hash:
                return True
            return slide1._serialized_elements == slide2._serialized_elements
        # Fallback to metadata comparison
        return slide1.title == slide2.title and slide1.layout_type == slide2.layout_type and slide1.element_count == slide2.element_count and slide1.notes == slide2.notes

    def _animation_key(self, animation: AnimationEffect) -> str:
        """Generate a unique key for an animation"""
        return f"{animation.slide_id}_{animation.element_id}_{animation.order}"

    def _animations_equal(self, anim1: AnimationEffect, anim2: AnimationEffect) -> bool:
        """Check if two animations are equal"""
        return anim1.effect_type == anim2.effect_type and anim1.trigger == anim2.trigger and abs(anim1.delay - anim2.delay) < 0.01 and abs(anim1.duration - anim2.duration) < 0.01

    def _transitions_equal(self, trans1: SlideTransition, trans2: SlideTransition) -> bool:
        """Check if two transitions are equal"""
        return trans1.transition_type == trans2.transition_type and abs(trans1.duration - trans2.duration) < 0.01 and trans1.direction == trans2.direction

    def generate_slide_screenshots(
        self,
        pptx_path: str,
        output_dir: Optional[str] = None,
        image_format: str = "jpg",
        density: int = 400,
        conversion_mode: str = "online",
        resolution: tuple|None = None,
    ) -> List[SlideScreenshot]:
        """
        Generate screenshots for all slides in a PowerPoint presentation.

        Args:
            pptx_path: Path to the PowerPoint file
            output_dir: Directory to save screenshots (default: temp directory)
            image_format: Image format (jpg, png, etc.)
            density: Image density/DPI for conversion (higher = better quality)

        Returns:
            List of SlideScreenshot objects with paths to generated images
        """
        return self.screenshot_generator.generate_slide_screenshots(
            pptx_path=pptx_path,
            output_dir=output_dir,
            image_format=image_format,
            density=density,
            conversion_mode=conversion_mode,
            resolution=resolution
        )

    def load_slide_screenshots(
        self,
        zip_path: str,
        output_dir: Optional[str] = None,
        pptx_path: Optional[str] = None,
    ) -> List[SlideScreenshot]:
        """
        Load slide screenshots from a local ZIP file containing slide images.

        Args:
            zip_path: Path to the ZIP file containing slide images
            output_dir: Directory to extract and save screenshots (default: temp directory)
            pptx_path: Optional path to the PowerPoint presentation file for proper slide_id mapping

        Returns:
            List of SlideScreenshot objects with paths to extracted images
        """
        import tempfile

        # Create output directory if not specified
        if output_dir is None:
            output_dir = tempfile.mkdtemp(prefix="zip_screenshots_")
        else:
            os.makedirs(output_dir, exist_ok=True)

        # Validate ZIP file exists
        if not os.path.exists(zip_path):
            raise FileNotFoundError(f"ZIP file not found: {zip_path}")

        # Use the screenshot generator's ZIP conversion method
        return self.screenshot_generator._convert_zip_to_screenshots(zip_path, output_dir, pptx_path)

    def slide_xml_to_json(self, slide_xml: bytes) -> Dict[str, Any]:
        """
        Convert slide XML to JSON format with attributes as key-values using xmltodict.

        Args:
            slide_xml: Raw XML bytes of the slide

        Returns:
            Dict representing the XML structure with attributes as key-values
        """
        try:
            # Convert bytes to string if needed
            if isinstance(slide_xml, bytes):
                slide_xml_str = slide_xml.decode('utf-8')
            else:
                slide_xml_str = slide_xml

            # Use xmltodict to convert XML to dictionary
            # The xmltodict library automatically handles attributes by prefixing them with '@'
            result = xmltodict.parse(slide_xml_str)

            return result
        except Exception as e:
            raise ValueError(f"Error converting XML to JSON using xmltodict: {e}")

